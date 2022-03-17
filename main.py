# -*- coding: utf-8 -*-
"""
This goes through the files in a give folder and its subfolders
and extracts certain types of metadata
also, if the textPull parameter is set to True then this will extract
the text from excel files, PDFs, word files, and PPTs
and if formulas is set to True then this will pull formulas out of Excel file text, if you have also set textPull to true

you can also modify the wordList function if you want to search the text you've extracted for a particular set of words

possible ways to expand this:
    --let you put a create or modify date in and ignore everything before/after that -- this could improve runtime a lot
    --expand to other types of files: .csv, .xls, .xlsb, 
    --figure out how to get this to run independently on Windows
    --add in more data engineering to have this write out periodically/let you pause/otherwise deal with the fact that
    this just takes a really long time to run if you're extracting text
    --do other kinds of filtering
    --try to extract info from file names (assuming certain standard naming conventions)
    --can we organize by subject -- topic modelling!!! and copying files over to new folder structure
    --be cleaner about importing libraries so if you aren't able to import certain libraries, this will still run and just not do certain things

resolved:
    no length of string that's going to wind up being a limitation
"""

### os, time, and pandas are packages everyone is going to have
### some of these are less so, but I was able to install all of them on my gfe
### if you are not able to, you can comment them out and run this with textPull=False

import os
import pandas as pd
import time
from openpyxl import load_workbook
from docx import Document
from pptx import Presentation
import PyPDF2

def conditions(file):
    # inputs: filename
    # outputs: True if it doesn't have $ and is isn't thumbs.db, false otherwise
    if "$" in file:
        return(False)
    if "Thumbs.db" in file:
        return(False)
    else:
        return(True)

def getCSVText(file):
    # inputs: filename of CSV file
    # outputs: all of the text in that file
    data=pd.read_csv(file).dropna(how='all').dropna(axis='columns',how='all')
    text=data.to_string().strip().replace("  ", " ").replace("\n","")
    return(text)
    
def runThroughPulls(file):
    # inputs: file name of Excel file
    # outputs: file text if able to read it, "" otherwise
    lower=file.lower()
    text=""
    try:
        if ".xls" in lower:
            creator, modified, text=scrapeExcel(file)
        elif ".ppt" in lower:
            creator, modified, text=scrapePPT(file)
        elif ".pdf" in lower:
            creator, modified, text=scrapePDF(file)
        elif ".doc" in lower:
            creator, modified, text=scrapeWord(file)
        elif ".csv" in lower:
            creator, modified=("", "")
            text=getCSVText(file)
        else:
            creator, modified, text="", "", ""
    except:
        creator, modified, text="", "", ""
    return(creator, modified, text)
    
def getTime(file):
    # inputs: fileName
    # outputs: creation and modify time of that file, if able to pull them
    try:
        creationTime=time.ctime(os.path.getctime(file))
        modifyTime=time.ctime(os.path.getmtime(file))
        return(creationTime, modifyTime)
    except:
        return("", "")

def runWords(text, keywords):
    # inputs: a string and the keywords we want to find in that string
    # outputs: the subset of keywords we found
    keywordsFound=[]
    for substring in keywords:
        if substring in text:
            keywordsFound.append(substring)
    return( ",".join(keywordsFound))
        
def getFileList(folder, keywords, textPull):
    # inputs: the folder name and keywords we want to find in the files there
    # outputs a df containing metadata on each xlsx file and text data from that file
    allFiles=[]
    for dirpath, dirnames, filenames in os.walk(folder):
        for filename in filenames:  
                file=os.path.join(dirpath,filename)
                if conditions(file):   
                    if textPull:
                        creator, modified, text=runThroughPulls(file)
                    else:
                        creator, modified, text="", "", ""
                    creationTime, modifyTime= getTime(file)
                    extension=file.split(".")[-1]
                    row=dirpath, filename, text, extension, creator, modified, creationTime, modifyTime
                    allFiles.append(row)
    df=pd.DataFrame(allFiles)
    df.columns=["location", "file name", "file text", "extension", "creator", "modified by", "creation time", "modified time"]
    df['Found keyword']=df['file text'].astype(str).str.lower().astype(str).apply(runWords, args=[keywords])
    df=df.drop_duplicates()
    if textPull==False:
        df=df.drop(columns=["creator", "modified by", "file text", "Found keyword"])
    return(df)
    
def wordList():
    # outputs: a list of words we're going to look for in our text
    keywords=["model", "inputs", "assumptions", "outputs",  "actual","predicted", "attributes"]
    return(keywords)
    
def scrapeWord(file):
    #inputs: file name (word doc)
    #outputs: creator, modified by, text of file
    try:
        document= Document(file)
        creator, modifiedBy=getInfo(document)
        text=getText(document)
        return(creator, modifiedBy, text)
    except:
        return("", "", "")

def getInfo(fileContent):
    #input: fileContent of a document or ppt
    #output: creator, modified by
    try:
        creator=fileContent.core_properties.author
        modifiedBy=fileContent.core_properties.last_modified_by
        return(creator, modifiedBy)
    except:
        return("","")
        
def scrapePPT(file):
    #inputs: filename of a ppt file
    #outputs: creator, modified by, filetext
    try:
        ppt= Presentation(file)
        creator, modifiedBy=getInfo(ppt)
        text=getPPTText(ppt)
        return(creator, modifiedBy, text)
    except:
        return("", "", "")
        
def scrapeExcel(file):
    # inputs: filename
    # outputs: creator, modified by, text of file
    try:
        wb = load_workbook(file)
        creator=wb.properties.creator
        modifiedBy=wb.properties.lastModifiedBy
        text=getExcelText(wb)
        return(creator, modifiedBy, text)
    except:
        return("", "","")

def getExcelText(wb):
    #inputs: workbook (openpyxl)
    #outputs: string with text containing all the text from the workbook
    cell_rows=[]
    for sheet in wb.worksheets:
        for row_cells in  sheet.iter_rows():
            cells=[cell.value for cell in row_cells if cell.value is not None]
            cell_rows.append(cells)
    flat_list = [str(item) for sublist in cell_rows for item in sublist]
    return(" ".join(flat_list))
        
def scrapePDF(file):
    #inputs: filename of PDF
    #outputs: text from PDF with the caveat that PDFs are twitchy and might not get everything
    pdfFileObject = open(file, 'rb')
    pdfReader = PyPDF2.PdfFileReader(pdfFileObject)
    creator=pdfReader.getDocumentInfo().author
    pages= pdfReader.numPages
    string=""
    for num in range(0, pages):
        pageObject = pdfReader.getPage(num)
        string=string+pageObject.extractText()
    pdfFileObject.close()
    text=string.replace("\n","")
    return(creator, "", text)
    
def getPPTText(ppt):
    #inputs: filename of PPT
    #ouputs: text from PPT
    text=""
    for slide in ppt.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text=text+(shape.text)+" "
    return(text)

def getText(document):
    #inputs: a document in docx
    #output: the text from that document
    allText=""
    for para in document.paragraphs:
        allText=allText+" " + para.text.strip()
    return(allText.strip().replace("  ", " "))
    
def getFormulas(string, extension):
    #inputs: string of text, extension of file
    #outputs: if this is an xlsx, this will return a list of the formulas in the filetext
    if extension=="xlsx":
        theList=string.split(" ")
        formulas=[i for i in theList if i.startswith("=")]
        return(",".join(formulas))
    else:
        return("")
        
def main(folder=os.getcwd(), textPull=True, formulas=False, writeOut=True):
    # inputs: folder we want to search - for instance, the output of an os.getcwd(),
    #as well as whether you want to pull the text out of documents if possible and whether you want to search excel output for formulas
    # outputs: a dataframe with text data and metadata from the excel files there
    keywords= wordList()
    df=getFileList(folder, keywords, textPull)
    if formulas:
        df['formulas']=df.apply(lambda x: getFormulas(x['file text'], x['extension']), axis=1)
    if writeOut:
        df.to_excel("fileList.xlsx")
    return(df) 
    
def testFolder(folder):
    #inputs: this takes a string that is supposed to be a folder name
    #outputs: True if you can access this folder, False if you cannot
    response=False
    currentDir=os.getcwd()
    try:
        os.chdir()
        response=True
    except:
        print("The folder you have entered is not properly formatted or you do not have access to it.")
        print(f'An example of a folder is your current directory: {currentDir}.')
    os.chdir(currentDir)
    return(response)
    
def getInputs():
    # outputs: this prompts users for parameters, runs the main function, and then returns a df
    # or else tells them that the parameters are wrong and prompts to exit or re-run
    response=False
    while response==False:
        folder = input("Enter the folder you wish to search: ")
    while response==True 
    textPull=input("I want to extract text from the files you're searching? (True/False): ")
    
    
